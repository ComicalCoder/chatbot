import os
import PyPDF2
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

INPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data")
OUTPUT_DIR = INPUT_DIR

def parse_pdf(filepath):
    try:
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            text = "".join(page.extract_text() or "" for page in reader.pages)
        return text
    except Exception as e:
        print(f"Error parsing PDF {filepath}: {e}")
        return ""

def parse_docx(filepath):
    try:
        doc = Document(filepath)
        text = "\n".join(para.text for para in doc.paragraphs)
        return text
    except Exception as e:
        print(f"Error parsing DOCX {filepath}: {e}")
        return ""

def parse_pptx(filepath):
    try:
        prs = Presentation(filepath)
        text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        return text
    except Exception as e:
        print(f"Error parsing PPTX {filepath}: {e}")
        return ""

def parse_xlsx(filepath):
    try:
        wb = load_workbook(filepath, read_only=True)
        text = ""
        for sheet in wb:
            for row in sheet.rows:
                text += " ".join(str(cell.value or "") for cell in row) + "\n"
        return text
    except Exception as e:
        print(f"Error parsing XLSX {filepath}: {e}")
        return ""

def parse_sm2(filepath):
    # Placeholder: Update when SM2 format is clarified
    print(f"SM2 parsing not implemented for {filepath}. Awaiting format details.")
    return ""

def save_text(filename, text):
    if not text:
        return
    output_path = os.path.join(OUTPUT_DIR, f"{os.path.splitext(filename)[0]}.txt")
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)

def main():
    for filename in os.listdir(INPUT_DIR):
        filepath = os.path.join(INPUT_DIR, filename)
        if not os.path.isfile(filepath):
            continue
        ext = os.path.splitext(filename)[1].lower()
        parsers = {
            ".pdf": parse_pdf,
            ".docx": parse_docx,
            ".pptx": parse_pptx,
            ".xlsx": parse_xlsx,
            ".sm2": parse_sm2
        }
        parser = parsers.get(ext)
        if parser:
            print(f"Parsing: {filename}")
            text = parser(filepath)
            save_text(filename, text)

if __name__ == "__main__":
    main()